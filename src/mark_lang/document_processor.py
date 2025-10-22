"""Document processing module for extracting text from various file formats."""
from __future__ import annotations

from pathlib import Path
from typing import Dict, List

import docx
import pdfplumber
from openpyxl import load_workbook
from pptx import Presentation


class DocumentProcessor:
    """Extract text content from various document formats."""

    def process_file(self, file_path: str | Path) -> Dict[str, str | List[str]]:
        """
        Process a document file and extract text content.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Dictionary with metadata and extracted text
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        extension = file_path.suffix.lower()
        
        if extension == ".pptx":
            return self._process_pptx(file_path)
        elif extension == ".docx":
            return self._process_docx(file_path)
        elif extension == ".pdf":
            return self._process_pdf(file_path)
        elif extension in [".xlsx", ".xls"]:
            return self._process_excel(file_path)
        else:
            raise ValueError(f"Unsupported file format: {extension}")
    
    def _process_pptx(self, file_path: Path) -> Dict[str, str | List[str]]:
        """Extract text from PowerPoint presentations."""
        prs = Presentation(file_path)
        slides_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                # Extract text from tables
                if hasattr(shape, "table"):
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            slide_text.append(row_text)
            
            if slide_text:
                slides_content.append({
                    "slide": slide_num,
                    "content": "\n".join(slide_text)
                })
        
        return {
            "file_name": file_path.name,
            "file_type": "pptx",
            "slides": slides_content,
            "full_text": "\n\n".join(slide["content"] for slide in slides_content)
        }
    
    def _process_docx(self, file_path: Path) -> Dict[str, str | List[str]]:
        """Extract text from Word documents."""
        doc = docx.Document(file_path)
        paragraphs = []
        tables_content = []
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text.strip())
        
        # Extract tables
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    table_text.append(row_text)
            if table_text:
                tables_content.append("\n".join(table_text))
        
        full_text = "\n".join(paragraphs)
        if tables_content:
            full_text += "\n\nTables:\n" + "\n\n".join(tables_content)
        
        return {
            "file_name": file_path.name,
            "file_type": "docx",
            "paragraphs": paragraphs,
            "tables": tables_content,
            "full_text": full_text
        }
    
    def _process_pdf(self, file_path: Path) -> Dict[str, str | List[str]]:
        """Extract text from PDF files."""
        pages_content = []
        
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    pages_content.append({
                        "page": page_num,
                        "content": text.strip()
                    })
                
                # Extract tables
                tables = page.extract_tables()
                if tables:
                    for table in tables:
                        table_text = "\n".join(
                            " | ".join(str(cell) for cell in row if cell)
                            for row in table if row
                        )
                        if table_text:
                            pages_content.append({
                                "page": page_num,
                                "content": f"[TABLE]\n{table_text}"
                            })
        
        return {
            "file_name": file_path.name,
            "file_type": "pdf",
            "pages": pages_content,
            "full_text": "\n\n".join(page["content"] for page in pages_content)
        }
    
    def _process_excel(self, file_path: Path) -> Dict[str, str | List[str]]:
        """Extract text from Excel files."""
        wb = load_workbook(file_path, read_only=True, data_only=True)
        sheets_content = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows_text = []
            
            for row in sheet.iter_rows(values_only=True):
                # Filter out empty cells and convert to string
                row_values = [str(cell) for cell in row if cell is not None and str(cell).strip()]
                if row_values:
                    rows_text.append(" | ".join(row_values))
            
            if rows_text:
                sheets_content.append({
                    "sheet": sheet_name,
                    "content": "\n".join(rows_text)
                })
        
        wb.close()
        
        return {
            "file_name": file_path.name,
            "file_type": "xlsx",
            "sheets": sheets_content,
            "full_text": "\n\n".join(f"[Sheet: {s['sheet']}]\n{s['content']}" for s in sheets_content)
        }
    
    def process_directory(self, directory_path: str | Path, extensions: List[str] | None = None) -> List[Dict]:
        """
        Process all documents in a directory.
        
        Args:
            directory_path: Path to the directory
            extensions: List of file extensions to process (default: all supported)
            
        Returns:
            List of extracted document contents
        """
        directory_path = Path(directory_path)
        
        if extensions is None:
            extensions = [".pptx", ".docx", ".pdf", ".xlsx", ".xls"]
        
        results = []
        for file_path in directory_path.iterdir():
            if file_path.is_file() and file_path.suffix.lower() in extensions:
                try:
                    content = self.process_file(file_path)
                    results.append(content)
                except Exception as e:
                    print(f"Error processing {file_path.name}: {e}")
        
        return results
